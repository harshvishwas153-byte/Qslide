import base64
import os
import tempfile
import zipfile
from dataclasses import dataclass

from pptx import Presentation
from pypdf import PdfReader
from pypdf import PdfWriter


@dataclass(frozen=True)
class CompressionResult:
    original_size: int
    final_size: int
    was_compressed: bool
    under_target: bool


_TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
)
_TINY_JPEG = base64.b64decode(
    "/9j/4AAQSkZJRgABAQAAAQABAAD/2wBDAP//////////////////////////////////////////////////////////////////////////////////////"
    "2wBDAf//////////////////////////////////////////////////////////////////////////////////////"
    "wAARCAABAAEDASIAAhEBAxEB/8QAFQABAQAAAAAAAAAAAAAAAAAAAAX/xAAUEAEAAAAAAAAAAAAAAAAAAAAA/9oADAMBAAIQAxAAAAH/xAAUEAEAAAAAAAAAAAAAAAAAAAAA/9oACAEBAAEFAqf/xAAUEQEAAAAAAAAAAAAAAAAAAAAA/9oACAEDAQE/ASP/xAAUEQEAAAAAAAAAAAAAAAAAAAAA/9oACAECAQE/ASP/xAAUEAEAAAAAAAAAAAAAAAAAAAAA/9oACAEBAAY/Aqf/xAAUEAEAAAAAAAAAAAAAAAAAAAAA/9oACAEBAAE/IX//2gAMAwEAAgADAAAAEP/EFBQRAQAAAAAAAAAAAAAAAAAAABD/2gAIAQMBAT8QH//EFBQRAQAAAAAAAAAAAAAAAAAAABD/2gAIAQIBAT8QH//EFBABAQAAAAAAAAAAAAAAAAAAABD/2gAIAQEAAT8QH//Z"
)
_TINY_GIF = base64.b64decode("R0lGODlhAQABAIAAAAAAAP///ywAAAAAAQABAAACAUwAOw==")
_TINY_BMP = base64.b64decode("Qk1GAAAAAAAAADYAAAAoAAAAAQAAAAEAAAABABgAAAAAABAAAADEDgAAxA4AAAAAAAAAAAAA////AA==")
_TINY_SVG = b'<svg xmlns="http://www.w3.org/2000/svg" width="1" height="1"/>'

_PPTX_HEAVY_PREFIXES = (
    "ppt/media/",
    "ppt/embeddings/",
    "ppt/activeX/",
)
_PPTX_PLACEHOLDERS = {
    ".png": _TINY_PNG,
    ".apng": _TINY_PNG,
    ".jpg": _TINY_JPEG,
    ".jpeg": _TINY_JPEG,
    ".jpe": _TINY_JPEG,
    ".gif": _TINY_GIF,
    ".bmp": _TINY_BMP,
    ".svg": _TINY_SVG,
}


def _same_dir_temp_path(file_path):
    fd, temp_path = tempfile.mkstemp(
        prefix=".compressed_",
        suffix=os.path.splitext(file_path)[1],
        dir=os.path.dirname(file_path) or None,
    )
    os.close(fd)
    return temp_path


def _compressed_pptx_part(name, data):
    lower_name = name.lower()
    if not lower_name.startswith(_PPTX_HEAVY_PREFIXES):
        return data

    extension = os.path.splitext(lower_name)[1]
    return _PPTX_PLACEHOLDERS.get(extension, b"")


def _compress_pptx_in_place(file_path):
    temp_path = _same_dir_temp_path(file_path)

    try:
        with zipfile.ZipFile(file_path, "r") as source, zipfile.ZipFile(
            temp_path,
            "w",
            compression=zipfile.ZIP_DEFLATED,
            compresslevel=9,
        ) as target:
            for source_info in source.infolist():
                target_info = zipfile.ZipInfo(source_info.filename, source_info.date_time)
                target_info.compress_type = zipfile.ZIP_DEFLATED
                target_info.external_attr = source_info.external_attr
                target_info.comment = source_info.comment

                data = b"" if source_info.is_dir() else source.read(source_info.filename)
                target.writestr(
                    target_info,
                    _compressed_pptx_part(source_info.filename, data),
                    compresslevel=9,
                )
        return temp_path
    except Exception:
        try:
            os.remove(temp_path)
        except OSError:
            pass
        raise


def _compress_pdf_in_place(file_path):
    temp_path = _same_dir_temp_path(file_path)

    try:
        reader = PdfReader(file_path)
        writer = PdfWriter()

        for page in reader.pages:
            page.compress_content_streams()
            writer.add_page(page)

        writer.remove_images()
        writer.compress_identical_objects()

        with open(temp_path, "wb") as output:
            writer.write(output)
        return temp_path
    except Exception:
        try:
            os.remove(temp_path)
        except OSError:
            pass
        raise


def compress_upload_for_processing(file_path, target_bytes):
    original_size = os.path.getsize(file_path)
    lower_path = file_path.lower()

    if original_size <= target_bytes:
        return CompressionResult(
            original_size=original_size,
            final_size=original_size,
            was_compressed=False,
            under_target=True,
        )

    temp_path = None
    try:
        if lower_path.endswith(".pptx"):
            temp_path = _compress_pptx_in_place(file_path)
        elif lower_path.endswith(".pdf"):
            temp_path = _compress_pdf_in_place(file_path)
        else:
            return CompressionResult(
                original_size=original_size,
                final_size=original_size,
                was_compressed=False,
                under_target=original_size <= target_bytes,
            )

        compressed_size = os.path.getsize(temp_path)
        if compressed_size < original_size:
            os.replace(temp_path, file_path)
            final_size = compressed_size
            was_compressed = True
        else:
            os.remove(temp_path)
            final_size = original_size
            was_compressed = False

        return CompressionResult(
            original_size=original_size,
            final_size=final_size,
            was_compressed=was_compressed,
            under_target=final_size <= target_bytes,
        )
    except Exception:
        if temp_path:
            try:
                os.remove(temp_path)
            except OSError:
                pass
        return CompressionResult(
            original_size=original_size,
            final_size=os.path.getsize(file_path),
            was_compressed=False,
            under_target=os.path.getsize(file_path) <= target_bytes,
        )

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        raise ValueError(
            "This PowerPoint file could not be read. Please upload a valid .pptx file, "
            "or export the presentation as PDF and upload the PDF."
        ) from exc

    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text

def extract_text_from_pdf(file_path):
    try:
        reader = PdfReader(file_path)
    except Exception as exc:
        raise ValueError(
            "This PDF file could not be read. Please upload a valid text-based PDF."
        ) from exc

    pages = []

    for page in reader.pages:
        pages.append(page.extract_text() or "")

    return "\n".join(pages)

def extract_text_from_file(file_path):
    lower_path = file_path.lower()

    if lower_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    if lower_path.endswith(".ppt"):
        raise ValueError(
            "Old .ppt files are not supported. Please save the presentation as .pptx "
            "or export it as PDF, then upload again."
        )
    if lower_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)

    raise ValueError("Unsupported file type. Please upload a PPTX or PDF file.")
